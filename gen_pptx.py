from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

NAVY   = RGBColor(0,  51, 102)
BLUE   = RGBColor(0,  90, 160)
CYAN   = RGBColor(0, 180, 220)
WHITE  = RGBColor(255, 255, 255)
LGRAY  = RGBColor(230, 235, 245)
DGRAY  = RGBColor(80,  80,  80)
GREEN  = RGBColor(0,  160,  80)
ORANGE = RGBColor(220, 100,  0)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

IMG      = "C:/CapstoneProject2/riscv_vpu/report/image/"
IMG_FPGA = "C:/CapstoneProject2/riscv_vpu/fpga/verify_out/"


def rect(slide, l, t, w, h, fill=NAVY):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return tb

def header(slide, title):
    rect(slide, 0, 0, 13.33, 7.5, LGRAY)
    rect(slide, 0, 0, 13.33, 1.0, NAVY)
    rect(slide, 0, 1.0, 13.33, 0.04, CYAN)
    txt(slide, title, 0.3, 0.1, 12.5, 0.8, size=24, bold=True, color=WHITE)

def bullets(slide, items, x=0.4, y=1.3, w=12.5):
    for (lvl, t) in items:
        prefix = "  • " if lvl == 0 else "      – "
        c = NAVY if lvl == 0 else DGRAY
        sz = 18 if lvl == 0 else 15
        txt(slide, prefix + t, x, y, w, 0.42, size=sz, bold=(lvl==0), color=c)
        y += 0.45

def pic(slide, path, l, t, w, h):
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))

# ── Slide 1: Title ────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
rect(s, 0, 0, 13.33, 7.5, NAVY)
rect(s, 0, 0, 13.33, 0.08, CYAN)
rect(s, 0, 7.42, 13.33, 0.08, CYAN)
txt(s, "RISC-V Vector Processing Unit", 0.8, 1.6, 11.5, 1.2, size=38, bold=True,
    color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "Targeting ASIC Tapeout", 0.8, 2.8, 11.5, 0.8, size=30, bold=True,
    color=CYAN, align=PP_ALIGN.CENTER)
rect(s, 3.5, 3.7, 6.3, 0.06, CYAN)
txt(s, "rv32im_zicsr_zve32x_zvl128b  |  VLEN=128  |  Single Execution Lane",
    0.8, 3.9, 11.5, 0.7, size=18, color=LGRAY, align=PP_ALIGN.CENTER)
txt(s, "Capstone Project 2  ·  Pham Anh Khoi  ·  HCMUT  ·  June 2026",
    0.8, 5.2, 11.5, 0.6, size=16, color=RGBColor(160, 200, 240), align=PP_ALIGN.CENTER)

# ── Slide 2: Goals ────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
header(s, "Project Overview & Goals")
bullets(s, [
    (0, "Target ISA: rv32im_zicsr_zve32x_zvl128b"),
    (1, "VLEN=128 bits, SEW=8/16/32, LMUL=1/2/4/8 -- full Zve32x subset"),
    (0, "Design Philosophy: ASIC-ready RTL"),
    (1, "Synthesisable SystemVerilog -- no latches, sync active-low reset throughout"),
    (1, "Parameterised interfaces -- no hardcoded VLEN/SEW/LMUL bit-slices"),
    (0, "Application: image processing accelerator"),
    (1, "BT.601 RGB to Grayscale on 128x128 Lena image"),
    (1, "FPGA demo on DE10-Standard (Cyclone V) with VGA output"),
    (0, "Toolchain: ModelSim simulation, Quartus synthesis, GCC RISC-V assembler"),
])

# ── Slide 3: System Architecture ─────────────────────────────────────────────
s = prs.slides.add_slide(blank)
header(s, "System Architecture")
pic(s, IMG + "TOP_design.png", 0.3, 1.1, 12.7, 5.8)
txt(s, "riscv_vpu_top: 5-stage scalar core + VPU + DMEM/IMEM + UART on one clock domain",
    0.4, 6.9, 12.5, 0.45, size=13, color=DGRAY, align=PP_ALIGN.CENTER)

# ── Slide 4: Scalar Pipeline ──────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
header(s, "5-Stage Scalar Pipeline (RV32IM)")
bullets(s, [
    (0, "IF / ID / EX / MEM / WB with synchronous 1-cycle IMEM/DMEM"),
    (0, "Full forwarding: EX->MEM and MEM->WB bypass muxes"),
    (0, "Load-use hazard: stall PC + IF/ID, inject bubble into EX"),
    (0, "VPU dispatch: OP-V in EX drives vpu_insn_vld; back-pressure via vpu_ready"),
    (1, "vsetvl* result written back to scalar RF via cfg_done pulse"),
], x=0.4, y=1.3, w=8.5)
pic(s, IMG + "Top-level_block_diagram.png", 8.8, 1.2, 4.2, 5.8)

# ── Slide 5: VPU Architecture ─────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
header(s, "VPU Subsystem Architecture")
pic(s, IMG + "vpu_top_wrapper.png", 0.3, 1.1, 12.7, 5.8)
txt(s, "Decoder (32b->49b ctrl_bus) -> FIFO (depth 8) -> FSM -> 4xVRF + 4 lanes + VLSU + Reduction",
    0.4, 6.9, 12.5, 0.45, size=13, color=DGRAY, align=PP_ALIGN.CENTER)

# ── Slide 6: FSM ─────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
header(s, "VPU Execution FSM -- 9 States")
pic(s, IMG + "fsm_diagram.png", 0.3, 1.1, 12.7, 5.7)
txt(s, "ST_IDLE is the central dispatch hub; every path returns to IDLE within 1 extra cycle after completion",
    0.4, 6.9, 12.5, 0.45, size=13, color=DGRAY, align=PP_ALIGN.CENTER)

# ── Slide 7: Execution Lane ───────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
header(s, "Execution Lane (vproc_processor_lane x4)")
pic(s, IMG + "executing_lane.png", 0.3, 1.1, 12.7, 5.8)
txt(s, "7 parallel functional units: Adder | Multiplier | Shifter | Logic | Compare | MinMax | Reduction",
    0.4, 6.9, 12.5, 0.45, size=13, color=DGRAY, align=PP_ALIGN.CENTER)

# ── Slide 8: VLSU Waveform ────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
header(s, "Vector LSU -- Masked Store VSE8.v")
pic(s, IMG + "vlsu_masked_store_wave.png", 0.3, 1.1, 12.7, 5.6)
txt(s, "vl=12, SEW=8, vm=0: 3-word transfer, mem_be=4b0101 every word (even elements active, odd masked)",
    0.4, 6.8, 12.5, 0.5, size=13, color=DGRAY, align=PP_ALIGN.CENTER)

# ── Slide 9: Verification ─────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
header(s, "Verification Results")
stats = [
    ("Instruction\nRegression", "172 / 172", BLUE),
    ("Vector LSU\nTests", "75 / 75", GREEN),
    ("Lena 128x128\nPixels", "16384 / 16384", GREEN),
    ("Max Pixel\nError", "ZERO", ORANGE),
]
bw = 2.8
for i, (lbl, val, clr) in enumerate(stats):
    x = 0.6 + i * 3.1
    rect(s, x, 2.0, bw, 3.2, clr)
    txt(s, val,  x, 2.3, bw, 1.4, size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, lbl,  x, 3.9, bw, 0.9, size=15, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "All test cases PASS -- 172 directed tests cover ALU/MUL/SHIFT/COMPARE/REDUCTION/VLSU/LMUL/WIDENING",
    0.5, 5.5, 12.3, 0.5, size=14, color=DGRAY, align=PP_ALIGN.CENTER)

# ── Slide 10: FPGA Synthesis ──────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
header(s, "FPGA Synthesis -- Cyclone V DE10-Standard")
pic(s, IMG + "Fmax.png", 0.3, 1.1, 12.7, 5.7)
txt(s, "Fast-corner: 47.18 MHz (+2.247 ns slack)  |  Slow-corner: -4.436 ns (scalar ALU carry chain bottleneck)",
    0.4, 6.9, 12.5, 0.45, size=13, color=DGRAY, align=PP_ALIGN.CENTER)

# ── Slide 11: Fmax evolution ─────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
header(s, "Fmax Optimization Journey")
bullets(s, [
    (0, "Initial: 16.44 MHz -- combinational 64-stage adder tree in vproc_reduction"),
    (1, "Critical path: 64x 32-bit adder stages in one always_comb block"),
    (0, "After FIFO restructure: 32.5 MHz"),
    (1, "Broke 143-node feedback loop: push_valid -> fifo_full -> vpu_ready -> push_valid"),
    (0, "Final (serial reduction): 47.18 MHz  -- 2.87x improvement"),
    (1, "One 32-bit adder per cycle; latency cost: N_elem extra cycles per reduction"),
    (1, "ALM count: 13,998 -> 10,438 (25% area reduction)"),
    (0, "Remaining bottleneck: scalar ALU carry chain (~6 ns slow-corner)"),
    (1, "Fix path: split adder into 2 pipeline stages or carry-skip optimization"),
])

# ── Slide 12: Demo Results ────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
header(s, "FPGA VGA Demo -- Lena 128x128")
rect(s, 0, 1.04, 13.33, 6.46, LGRAY)
txt(s, "Original RGB  (SW[0] = 1)", 1.5, 1.2, 4.5, 0.5,
    size=17, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
txt(s, "Grayscale BT.601  (SW[0] = 0)", 7.3, 1.2, 4.5, 0.5,
    size=17, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
pic(s, IMG_FPGA + "mif_rgb.png",  1.5, 1.75, 4.5, 4.5)
pic(s, IMG_FPGA + "mif_gray.png", 7.3, 1.75, 4.5, 4.5)
txt(s, "16,384 / 16,384 pixels verified byte-exact  |  max_error = 0  |  Hardware confirmed on DE10-Standard",
    0.5, 6.5, 12.3, 0.6, size=14, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# ── Slide 13: Performance ─────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
header(s, "Performance Benchmarks")
bullets(s, [
    (0, "BT.601 Lena 128x128 (SEW=8, 1024 iterations)"),
    (1, "35853 clock cycles  ->  2.18 cycles/pixel"),
    (0, "Matrix multiply 4x4 (SEW=32)"),
    (1, "317 cycles -- outer-product accumulation, correct"),
    (0, "AXPY N=16 (a=3)"),
    (1, "221 cycles -- y[0]=4, y[4]=16, y[8]=28, y[12]=40 correct"),
    (0, "VPU speedup vs scalar baseline"),
    (1, "~10x for BT.601 inner loop (16 pixels per iteration)"),
    (1, "Bandwidth-bound: 1 DMEM byte/cycle matches 1 element/cycle at SEW=8"),
])

# ── Slide 14: Conclusion ──────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
header(s, "Conclusion & Future Work")
bullets(s, [
    (0, "Achievements"),
    (1, "Complete Zve32x VPU -- 172/172 regression PASS"),
    (1, "End-to-end system verification: UART, image processing, FPGA demo"),
    (1, "47.18 MHz, 10,438 ALMs on Cyclone V -- 3x frequency and 25% area improvement"),
    (1, "Hardware VGA display verified on DE10-Standard board"),
    (0, "Path to ASIC Tapeout"),
    (1, "Fix async reset in lsu.sv (#8) -- required for synthesis"),
    (1, "Lint pass for latches (#9) -- pre-synthesis gate"),
    (1, "Replace behavioral DMEM with SRAM macro (#15)"),
    (0, "Feature Extensions"),
    (1, "vslide1up/vslide1down (FIR filter, overlap-save)"),
    (1, "Formal verification and static timing closure at 50 MHz"),
])

out = "C:/CapstoneProject2/riscv_vpu/RISC-V_VPU_Presentation.pptx"
prs.save(out)
print(f"Saved {out}  ({len(prs.slides)} slides)")
