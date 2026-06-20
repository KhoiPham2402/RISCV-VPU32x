#!/usr/bin/env python3
"""generate_report.py — Build Design_Report.pptx from Weekly_Report_template.pptx"""

import copy
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

TEMPLATE = 'report/Weekly_Report_template.pptx'
OUTPUT   = 'report/Design_Report.pptx'

A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
R_NS = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'

# ── Slide copy / remove ──────────────────────────────────────────────────────

def copy_slide(prs, src_idx):
    src = prs.slides[src_idx]
    new_sl = prs.slides.add_slide(src.slide_layout)
    dst = new_sl.shapes._spTree
    for el in list(dst):
        dst.remove(el)
    for el in src.shapes._spTree:
        new_el = copy.deepcopy(el)
        for blip in new_el.iter(f'{{{A_NS}}}blip'):
            k = f'{{{R_NS}}}embed'
            rid = blip.get(k)
            if rid and rid in src.part.rels:
                rel = src.part.rels[rid]
                new_rid = new_sl.part.relate_to(rel._target, rel.reltype)
                blip.set(k, new_rid)
        dst.append(new_el)
    return new_sl


def remove_slide(prs, idx):
    R_ID = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
    lst = prs.slides._sldIdLst
    rId = lst[idx].get(R_ID)
    prs.part.drop_rel(rId)
    del lst[idx]


def remove_shape_by_name(slide, name_part):
    for sh in slide.shapes:
        if name_part in sh.name:
            sh._element.getparent().remove(sh._element)
            return


# ── XML helpers ──────────────────────────────────────────────────────────────

def make_p(text, sz, bold=False, bullet=False, lang='vi-VN', color=None):
    """Build one <a:p> element."""
    p = etree.Element(f'{{{A_NS}}}p')
    pPr = etree.SubElement(p, f'{{{A_NS}}}pPr')
    if bullet:
        pPr.set('marL', '342900'); pPr.set('indent', '-342900')
        bu = etree.SubElement(pPr, f'{{{A_NS}}}buChar')
        bu.set('char', '▶')
    else:
        pPr.set('marL', '0'); pPr.set('indent', '0')
        etree.SubElement(pPr, f'{{{A_NS}}}buNone')
    r = etree.SubElement(p, f'{{{A_NS}}}r')
    rPr = etree.SubElement(r, f'{{{A_NS}}}rPr')
    rPr.set('lang', lang)
    rPr.set('sz', str(sz))
    if bold:
        rPr.set('b', '1')
    rPr.set('dirty', '0')
    if color:
        sf = etree.SubElement(rPr, f'{{{A_NS}}}solidFill')
        sc = etree.SubElement(sf,  f'{{{A_NS}}}srgbClr')
        sc.set('val', color)
    etree.SubElement(r, f'{{{A_NS}}}t').text = text
    return p


def set_txbody(shape, items):
    """items: [(text, sz, bold, bullet)]  optional 5th element = lang."""
    tb = shape.text_frame._txBody
    for p in list(tb.findall(f'{{{A_NS}}}p')):
        tb.remove(p)
    for item in items:
        text, sz, bold, bullet = item[:4]
        lang  = item[4] if len(item) > 4 else 'vi-VN'
        color = item[5] if len(item) > 5 else None
        tb.append(make_p(text, sz, bold, bullet, lang, color))


def find_sh(slide, name_part):
    for s in slide.shapes:
        if name_part in s.name:
            return s
    return None


def set_ph(slide, idx, items):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            set_txbody(ph, items)
            return


# ── Section label (TextBox matching template chrome style) ───────────────────

def add_section_label(slide, text):
    # Same position as Google Shape;116 on slides 2-10
    # L=1183907  T=327486  W=11715750  H=553998  sz=36pt bold dark-blue
    tb = slide.shapes.add_textbox(Emu(1183907), Emu(327486), Emu(11715750), Emu(553998))
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.bold  = True
    run.font.size  = Pt(28)
    run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)


# ── Diagram placeholder box ──────────────────────────────────────────────────

def add_diag_box(slide, label, left_emu, top_emu, w_emu, h_emu):
    """Dashed blue outline box with centered label — user replaces with diagram."""
    from pptx.oxml.ns import qn as _qn
    # Add via XML: <p:sp> with prstGeom rect + no fill + dashed border
    sp_xml = f"""<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvSpPr>
    <p:cNvPr id="9901" name="DiagBox"/>
    <p:cNvSpPr txBox="1"/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="{left_emu}" y="{top_emu}"/>
             <a:ext cx="{w_emu}"  cy="{h_emu}"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:noFill/>
    <a:ln w="19050" cmpd="sng">
      <a:solidFill><a:srgbClr val="4172C4"/></a:solidFill>
      <a:prstDash val="dash"/>
    </a:ln>
  </p:spPr>
  <p:txBody>
    <a:bodyPr anchor="ctr"><a:spAutoFit/></a:bodyPr>
    <a:lstStyle/>
    <a:p>
      <a:pPr algn="ctr"><a:buNone/></a:pPr>
      <a:r>
        <a:rPr lang="vi-VN" sz="1600" i="1" dirty="0">
          <a:solidFill><a:srgbClr val="4172C4"/></a:solidFill>
        </a:rPr>
        <a:t>{label}</a:t>
      </a:r>
    </a:p>
  </p:txBody>
</p:sp>"""
    sp_el = etree.fromstring(sp_xml)
    slide.shapes._spTree.append(sp_el)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE CONTENT
# (section_label, title, content_items, diagram_label_or_None)
# content_items: (text, sz, bold, bullet)
# ═══════════════════════════════════════════════════════════════════════════════

SLIDES = [

    # ── VLSU 1 ──────────────────────────────────────────────────────────────
    ("VLSU Design",
     "Vector Load/Store Unit — Kiến Trúc",
     [
         ("Chức năng",                                                   2200, True,  False),
         ("Thực hiện VLE8/16/32.V (load) và VSE8/16/32.V (store) unit-stride.",
                                                                          1500, False, True),
         ("Truy cập DMEM độc lập — parallel với scalar port, không tranh chấp bus.",
                                                                          1500, False, True),
         ("Kết nối trực tiếp dmem_sync qua VLSU dedicated port (không qua TileLink).",
                                                                          1500, False, True),
         ("Byte-enable & Masked Store",                                   2200, True,  False),
         ("SEW=e8 → 1 byte/word;  SEW=e16 → 2 bytes;  SEW=e32 → 4 bytes per cycle.",
                                                                          1500, False, True),
         ("vm=0 (masked): byte-enable = 0 cho element bị mask (vm_i + v0_flat_i ports).",
                                                                          1500, False, True),
         ("last_word_be: ngăn ghi tail bytes ngoài VL ở word cuối cùng.",
                                                                          1500, False, True),
         ("Handshake: vlsu_ready=1 ngay cho store;  =1 sau 1 cycle cho load.",
                                                                          1500, False, True),
     ],
     "Chèn sơ đồ khối VLSU tại đây"),

    # ── VLSU 2 ──────────────────────────────────────────────────────────────
    ("VLSU Design",
     "VLSU — Interface Signals & Protocol",
     [
         ("Ports:  vproc_vec_lsu.sv  ↔  dmem_sync.sv",                  2200, True,  False),
         ("vlsu_mem_req    — pulse 1 cycle kích hoạt giao dịch",         1500, False, True),
         ("vlsu_mem_we     — 1 = store,  0 = load",                      1500, False, True),
         ("vlsu_mem_addr   — byte address [31:0]",                       1500, False, True),
         ("vlsu_mem_be     — byte enables [3:0]  (SEW + vm mask)",       1500, False, True),
         ("vlsu_mem_wdata  — store data [31:0]",                         1500, False, True),
         ("vlsu_mem_rdata  — load data [31:0]  (valid 1 cycle sau req)", 1500, False, True),
         ("vlsu_mem_ready  — write: 1 ngay;  read: 1 sau 1 cycle",      1500, False, True),
         ("Conflict Resolution",                                          2200, True,  False),
         ("Scalar vs VLSU ghi cùng word: VLSU thắng (if-block VLSU sau scalar trong dmem_sync).",
                                                                          1500, False, True),
         ("vrf_busy_r[31:0] scoreboard ngăn RAW hazard VRF giữa vector instructions liên tiếp.",
                                                                          1500, False, True),
     ],
     "Chèn waveform handshake VLSU tại đây"),

    # ── Core Interface 1 ─────────────────────────────────────────────────────
    ("Core Interface",
     "Scalar ↔ VPU Interface Signals",
     [
         ("Scalar  →  VPU   (dispatch từ EX stage)",                     2200, True,  False),
         ("vpu_insn_vld_o   — valid khi vector instruction ở EX stage",  1500, False, True),
         ("vpu_insn_o       — instruction word [31:0]",                  1500, False, True),
         ("vpu_rs1_data_o   — scalar rs1 sau EX/MEM/WB forwarding",     1500, False, True),
         ("vpu_rs2_data_o   — scalar rs2 sau EX/MEM/WB forwarding",     1500, False, True),
         ("VPU  →  Scalar   (feedback / stall / config)",                2200, True,  False),
         ("vpu_ready_i      — VPU sẵn sàng nhận lệnh mới (FIFO có chỗ)",1500, False, True),
         ("vpu_cfg_done_i   — vsetvl*/vsetvli đã commit, vl mới hợp lệ",1500, False, True),
         ("vpu_vl_remain_i  — giá trị vl mới [31:0] → scalar ghi vào rd",1500, False, True),
     ],
     "Chèn sơ đồ interface Scalar-VPU tại đây"),

    # ── Core Interface 2 ─────────────────────────────────────────────────────
    ("Core Interface",
     "Stall, Dispatch & Hazard Handling",
     [
         ("Pipeline Stall   (pipelined_vpu.sv)",                         2200, True,  False),
         ("vpu_stall = is_vector_exe && !vpu_ready_i",                   1500, False, True),
         ("Khi stall: enable = !vpu_stall cho tất cả pipeline registers (giữ nguyên state).",
                                                                          1500, False, True),
         ("vsetvl* writeback: vpu_cfg_done → rf_wren=1, ghi vpu_vl_remain vào rd_addr.",
                                                                          1500, False, True),
         ("Hazard Control   (vproc_system_wrapper.sv)",                  2200, True,  False),
         ("RAW:  vrf_busy_r[31:0] — set khi dispatch, clear khi WB commit.",
                                                                          1500, False, True),
         ("WAW (load):  load_waw_stall — block store nếu VLSU load chưa ghi VRF.",
                                                                          1500, False, True),
         ("Duplicate-push fix: push_valid gated by vpu_ready_i — tránh push FIFO 2 lần.",
                                                                          1500, False, True),
     ],
     "Chèn timing diagram pipeline stall tại đây"),

    # ── Memory ────────────────────────────────────────────────────────────────
    ("Memory",
     "Memory Architecture",
     [
         ("Harvard Architecture",                                         2200, True,  False),
         ("IMEM (imem_sync.sv): 8KB, sync 1-cycle read, $readmemh(\"imem.hex\") lúc elaboration.",
                                                                          1500, False, True),
         ("DMEM (dmem_sync.sv): 64KB = 16384 × 32-bit words, sync reset active-low.",
                                                                          1500, False, True),
         ("dmem_sync — Dual-Port",                                        2200, True,  False),
         ("Scalar port: SB byte-enable từ addr[1:0]+func3; SW/SH/SB phân lane; 1-cycle read.",
                                                                          1500, False, True),
         ("VLSU port:  vlsu_req/we/addr/be/wdata/rdata/ready; write ngay; read +1 cycle.",
                                                                          1500, False, True),
         ("Conflict: VLSU thắng khi cả 2 ghi cùng word (if-guard trong always_ff).",
                                                                          1500, False, True),
         ("Memory Map — Lena 128×128 Benchmark",                         2200, True,  False),
         ("0x0000–0x3FFF   R channel  (16384 bytes)  ← nhận qua UART",  1500, False, True),
         ("0x4000–0x7FFF   G channel  (16384 bytes)  ← nhận qua UART",  1500, False, True),
         ("0x8000–0xBFFF   B channel  (16384 bytes)  ← nhận qua UART",  1500, False, True),
         ("0xC000–0xFFFF   Y output   (16384 bytes)  ← ghi bởi VPU VSE8.V",
                                                                          1500, False, True),
     ],
     "Chèn sơ đồ memory hierarchy tại đây"),

    # ── UART 1 ────────────────────────────────────────────────────────────────
    ("UART Integration",
     "UART via TileLink-UL Bus",
     [
         ("TileLink-UL Protocol   (rtl_trial/bus/tl_pkg.sv)",            2200, True,  False),
         ("Channel A (tl_a_t):  opcode | address[31:0] | mask[3:0] | data[31:0] | valid",
                                                                          1500, False, True),
         ("Channel D (tl_d_t):  opcode | data[31:0] | valid   (slave → master)",
                                                                          1500, False, True),
         ("GET (read), PUT_PARTIAL (write); D: ACCESS_ACK_DATA / ACCESS_ACK.",
                                                                          1500, False, True),
         ("UART Register Map   (uart.sv  @  0xFF000000)",                2200, True,  False),
         ("+0x00  CTRL   [0]=TX_EN,  [1]=RX_EN              default 1/1 tại reset",
                                                                          1500, False, True),
         ("+0x04  STAT   [0]=TX_BUSY, [1]=RX_VALID          RX FIFO non-empty flag",
                                                                          1500, False, True),
         ("+0x08  TXDAT  write [7:0] → push TX FIFO         triggers TX engine",
                                                                          1500, False, True),
         ("+0x0C  RXDAT  read  [7:0] → dequeue RX FIFO      side-effect on TL GET",
                                                                          1500, False, True),
         ("+0x10  BAUD   divisor = CLK_FREQ / (16 × BAUD_RATE) − 1",   1500, False, True),
         ("8-entry FIFO mỗi chiều; 16× oversampling; frame 8N1.",        1500, False, True),
     ],
     "Chèn sơ đồ TileLink A/D channel tại đây"),

    # ── UART 2 ────────────────────────────────────────────────────────────────
    ("UART Integration",
     "riscv_vpu_top_v4 — System Integration",
     [
         ("Address Decode  (inline top-level, không crossbar)",           2200, True,  False),
         ("uart_sel = (s_dmem_addr[31:8] == 24'hFF0000)",                1500, False, True),
         ("dmem_we = s_dmem_we & ~uart_sel   →  UART store không vào DMEM",
                                                                          1500, False, True),
         ("dmem_re = s_dmem_re & ~uart_sel   →  UART load không kích DMEM read",
                                                                          1500, False, True),
         ("1-Cycle Read Latency Alignment",                               2200, True,  False),
         ("UART TL-D trả lời combinatorial; DMEM có 1-cycle registered output.",
                                                                          1500, False, True),
         ("uart_rd_pending_r & uart_rdata_r registered 1 cycle để match latency.",
                                                                          1500, False, True),
         ("s_dmem_rdata = uart_rd_pending_r ? uart_rdata_r : dmem_rdata",1500, False, True),
         ("Firmware — sw/uart_lena.S  (53 instructions)",                2200, True,  False),
         ("Receive 3 × 16384 bytes R/G/B qua UART polling STAT[1]=RX_VALID.",
                                                                          1500, False, True),
         ("1024 iter VPU: vle8.v×3 + vmulhu.vx×3 + vadd.vv×2 + vse8.v (16 px/iter).",
                                                                          1500, False, True),
         ("Gửi ACK 0xAA;  total ~3.93 M cycles UART + VPU processing.",  1500, False, True),
     ],
     "Chèn block diagram riscv_vpu_top_v4 tại đây"),

    # ── Verification ──────────────────────────────────────────────────────────
    ("Verification",
     "Simulation & Kết Quả",
     [
         ("Tests đã PASS — VPU Core",                                    2200, True,  False),
         ("172/172  tb_vproc_all_instr: ALU VV/VX/VI, MUL, SHIFT, COMPARE, MINMAX, WIDENING, REDUCTION.",
                                                                          1500, False, True),
         ("75/75    tb_vproc_vlsu: e8/e16/e32, VL ngẫu nhiên, masked store.",
                                                                          1500, False, True),
         ("Tests đã PASS — Integration",                                  2200, True,  False),
         ("PASS    tb_lena_gray: 128×128 BT.601, 16384/16384 pixel đúng, max_err=0, 34828 cycles.",
                                                                          1500, False, True),
         ("PASS    tb_riscv_vpu_top: bench_imgproc, bench_matmul 4×4 int, bench_axpy N=16.",
                                                                          1500, False, True),
         ("Pending — UART Integration",                                   2200, True,  False),
         ("tb_uart_lena:     smoke test 16 pixel  R=G=B=128 → expected Y=127.",
                                                                          1500, False, True),
         ("tb_uart_lena_img: full 128×128 Lena qua UART → DMEM dump → reconstruct PNG.",
                                                                          1500, False, True),
         ("Chạy: vsim -c -do run_uart_sim.do  rồi  run_uart_lena_img_sim.do",
                                                                          1400, False, True),
     ],
     None),
]


# ═══════════════════════════════════════════════════════════════════════════════
# Build presentation
# ═══════════════════════════════════════════════════════════════════════════════

prs = Presentation(TEMPLATE)

# ── 1. Title slide (slide 0) ─────────────────────────────────────────────────
s0 = prs.slides[0]
for sh in s0.shapes:
    if sh.name == 'TextBox 5' and sh.has_text_frame:
        set_txbody(sh, [
            ("ĐỒ ÁN CHUYÊN NGÀNH 2",              2000, True,  False),
            ("BÁO CÁO THIẾT KẾ HỆ THỐNG",         2400, True,  False),
            ("RISCV + VECTOR PROCESSING UNIT",     2000, True,  False),
            ("VLSU  ·  Core Interface  ·  Memory  ·  UART", 1600, False, False),
        ])
        break

# ── 2. Contents slide (slide 1) ──────────────────────────────────────────────
s1 = prs.slides[1]
for sh in s1.shapes:
    if sh.name == 'Title 5' and sh.has_text_frame:
        set_txbody(sh, [("Nội Dung Báo Cáo", 3600, True, False)])
        break

for sh in s1.shapes:
    if sh.shape_type == 19:   # TABLE
        tbl = sh.table
        rows_data = [
            ("01", "VLSU Design",
             "Kiến trúc VLSU · interface dmem_sync · byte-enable · masked store"),
            ("02", "Core Interface",
             "Scalar ↔ VPU signals · stall protocol · RAW/WAW hazard"),
            ("03", "Memory Architecture",
             "imem_sync · dmem_sync dual-port · Harvard layout · memory map 64KB"),
            ("04", "UART Integration",
             "TileLink-UL · register map 0xFF000000 · address decode · firmware"),
        ]
        for r_i, (num, topic, desc) in enumerate(rows_data):
            row_idx = r_i + 1   # skip header row
            if row_idx >= len(tbl.rows):
                continue
            row = tbl.rows[row_idx]
            for c_i, text in enumerate([num, topic, desc]):
                if c_i >= len(tbl.columns):
                    continue
                cell = row.cells[c_i]
                cell.text_frame.clear()
                run = cell.text_frame.paragraphs[0].add_run()
                run.text = text
        break

# ── 3. Create content slides by copying slide 11 (full-width content) ────────
#   Slide 11 layout: Title (idx=0) + Content (idx=1, W≈8.27 in) + logo + line
BASE = 11

for section, title, content, diag_label in SLIDES:
    sl = copy_slide(prs, BASE)

    # Remove slide-11-specific pictures (keep logo=Picture 3, line, Picture 6)
    remove_shape_by_name(sl, 'Picture 15')
    remove_shape_by_name(sl, 'Picture 5')

    # Section label (top bar, same position as original slides 2-10)
    add_section_label(sl, section)

    # Title placeholder (idx=0)
    set_ph(sl, 0, [(title, 2800, True, False)])

    # Content placeholder (idx=1)
    # Reserve bottom rows if diagram needed
    if diag_label:
        # Keep top ~55% for text, bottom 45% will have diagram box
        # Content placeholder height = 4351338 EMU; top at 1325563
        # We'll shrink it by adjusting the XML xfrm height
        for ph in sl.placeholders:
            if ph.placeholder_format.idx == 1:
                set_txbody(ph, content)
                # Shrink height to leave room for diagram box below
                xfrm = ph._element.spPr.find(f'{{{A_NS}}}xfrm')
                if xfrm is not None:
                    ext = xfrm.find(f'{{{A_NS}}}ext')
                    if ext is not None:
                        ext.set('cy', str(2400000))   # ~1.9 in, leaves room below
                break

        # Diagram box: below content, full width area
        # T = 1325563 + 2400000 + 100000 = 3825563
        add_diag_box(sl, diag_label,
                     left_emu=838199, top_emu=3825563,
                     w_emu=10515600,  h_emu=2500000)
    else:
        set_ph(sl, 1, content)

# ── 4. Delete original slides 2-13 (12 slides) ───────────────────────────────
for _ in range(12):
    remove_slide(prs, 2)

prs.save(OUTPUT)
print(f'Saved: {OUTPUT}  ({len(prs.slides)} slides)')
