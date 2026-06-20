"""
gen_reports.py
Generates:
  1. report/Weekly_Report_Final.pptx  -- weekly progress report following HCMUT template style
  2. RISC-V_VPU_Presentation.pptx     -- clean summary presentation (content-first)

Run: python gen_reports.py
"""
import copy, io, os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

TEMPLATE  = r"C:\CapstoneProject2\riscv_vpu\report\Weekly_Report_template.pptx"
OUT_WR    = r"C:\CapstoneProject2\riscv_vpu\report\Weekly_Report_Final.pptx"
OUT_PPTX  = r"C:\CapstoneProject2\riscv_vpu\RISC-V_VPU_Presentation.pptx"

IMG  = r"C:\CapstoneProject2\riscv_vpu\report\image\\"
IMGF = r"C:\CapstoneProject2\riscv_vpu\fpga\verify_out\\"

# ─── helpers ──────────────────────────────────────────────────────────────────
def add_txt(slide, text, l, t, w, h, size=16, bold=False,
            color=RGBColor(30,30,30), align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.italic = italic
    return tb

def add_rect(slide, l, t, w, h, fill):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s

def add_pic(slide, path, l, t, w, h):
    if os.path.exists(path):
        return slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))

# ─── copy logo + divider from template slide ──────────────────────────────────
def copy_logo_header(src_slide, dst_slide):
    """Copy the BK logo picture and horizontal line from template slide to dst_slide."""
    for shape in src_slide.shapes:
        if shape.name in ("Picture 3",) or (shape.shape_type == 13 and
              abs(shape.left) < 200000 and shape.top < 1800000 and shape.width < 1500000):
            sp_elem = copy.deepcopy(shape.element)
            dst_slide.shapes._spTree.append(sp_elem)
        elif shape.shape_type == 9 and abs(shape.top - Inches(1.1)) < Inches(0.05):
            sp_elem = copy.deepcopy(shape.element)
            dst_slide.shapes._spTree.append(sp_elem)

NAVY  = RGBColor(0, 51, 102)
BLUE  = RGBColor(31, 73, 125)
WHITE = RGBColor(255, 255, 255)
LGRAY = RGBColor(242, 242, 242)
DGRAY = RGBColor(89, 89, 89)
GREEN = RGBColor(0, 128, 64)
RED   = RGBColor(192, 0, 0)
CYAN  = RGBColor(0, 176, 240)

# ════════════════════════════════════════════════════════════════════════════════
# PART 1 — WEEKLY REPORT (HCMUT template style)
# ════════════════════════════════════════════════════════════════════════════════
tpl = Presentation(TEMPLATE)
wr  = Presentation(TEMPLATE)

# Keep only slide 1 (title page) — remove the rest
slide_ids = list(wr.slides._sldIdLst)
for sid in slide_ids[1:]:
    rId = sid.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    if rId:
        try:
            wr.part.drop_rel(rId)
        except Exception:
            pass
    wr.slides._sldIdLst.remove(sid)

# Reference slide for logo/header copy
ref_slide = tpl.slides[1]   # slide 2 has logo + line

def wr_slide(title_top, section_label, content_fn):
    """Add a content slide to Weekly Report."""
    layout = tpl.slide_layouts[1]  # Title and Content
    slide = wr.slides.add_slide(layout)
    # Clear default placeholder content
    for ph in slide.placeholders:
        ph.text = ""
    copy_logo_header(ref_slide, slide)
    # Section label (top-left after logo)
    add_txt(slide, section_label, 1.3, 0.28, 11.5, 0.55,
            size=18, bold=True, color=NAVY)
    # Thin accent line
    add_rect(slide, 1.35, 0.92, 11.2, 0.04, NAVY)
    content_fn(slide)
    return slide

def bullets(slide, items, x=0.4, y=1.2, w=12.5, line_h=0.45):
    for (lvl, text) in items:
        c = NAVY if lvl == 0 else DGRAY
        sz = 17 if lvl == 0 else 14
        prefix = "◆  " if lvl == 0 else "     •  "
        add_txt(slide, prefix + text, x, y, w, line_h + 0.05,
                size=sz, bold=(lvl == 0), color=c)
        y += line_h + (0.05 if lvl == 0 else 0)
    return y

# ── Slide 1: Title — already in template, just update text ───────────────────
title_slide = wr.slides[0]
for shape in title_slide.shapes:
    if shape.name == 'TextBox 5' and shape.has_text_frame:
        tf = shape.text_frame
        tf.paragraphs[0].runs[0].text = "ĐỒ ÁN CHUYÊN NGÀNH 2"
        if len(tf.paragraphs) > 1:
            tf.paragraphs[1].runs[0].text = "BÁO CÁO CUỐI KỲ"
        if len(tf.paragraphs) > 2:
            tf.paragraphs[2].runs[0].text = "RISCV - VECTOR PROCESSING UNIT"

# ── Slide 2: Outline ──────────────────────────────────────────────────────────
def s2(slide):
    chapters = [
        "1. Giới thiệu & Mục tiêu đề tài",
        "2. Kiến trúc hệ thống tổng quan",
        "3. Scalar RISC-V Core (5-stage pipeline)",
        "4. Vector Processing Unit (VPU)",
        "   4.1 Decoder & ctrl_bus 49-bit",
        "   4.2 Instruction FIFO & FSM 9-state",
        "   4.3 Execution Lane & Functional Units",
        "   4.4 Vector LSU (VLE/VSE masked)",
        "   4.5 Reduction Unit (serial accumulator)",
        "5. Bộ nhớ & UART",
        "6. Phần mềm & Firmware",
        "7. Kiểm tra & Đánh giá",
        "8. Kết quả FPGA & Demo",
        "9. Kết luận & Hướng phát triển",
    ]
    y = 1.2
    for c in chapters:
        indent = c.startswith("   ")
        sz = 14 if indent else 16
        clr = DGRAY if indent else NAVY
        add_txt(slide, c, 0.6, y, 12.0, 0.38, size=sz, color=clr, bold=not indent)
        y += 0.38
wr_slide("MỤC LỤC", "Nội dung báo cáo", s2)

# ── Slide 3: Introduction ─────────────────────────────────────────────────────
def s3(slide):
    bullets(slide, [
        (0, "Mục tiêu: Thiết kế VPU tích hợp với RISC-V scalar core, hướng đến tapeout ASIC"),
        (1, "ISA: rv32im_zicsr_zve32x_zvl128b — bao gồm Zve32x subset đầy đủ"),
        (1, "VLEN = 128 bits, SEW = 8/16/32 bit, LMUL = 1/2/4/8"),
        (0, "Ứng dụng minh họa: Xử lý ảnh thời gian thực"),
        (1, "BT.601 RGB → Grayscale trên ảnh Lena 128×128"),
        (1, "Demo FPGA trên DE10-Standard (Cyclone V) với VGA output"),
        (0, "Quy tắc thiết kế RTL (ASIC-ready)"),
        (1, "Không dùng construct chỉ dành cho simulation trong /rtl/"),
        (1, "Synchronous active-low reset xuyên suốt toàn bộ hierarchy"),
        (1, "Parameterised interfaces — không hardcode VLEN/SEW/LMUL"),
    ])
wr_slide("CHƯƠNG 1", "Giới thiệu & Mục tiêu", s3)

# ── Slide 4: System Architecture ─────────────────────────────────────────────
def s4(slide):
    add_txt(slide, "Kiến trúc hệ thống tổng quan", 1.3, 1.15, 11.0, 0.5,
            size=16, bold=True, color=NAVY)
    add_pic(slide, IMG + "TOP_design.png", 0.5, 1.6, 12.3, 5.5)
wr_slide("CHƯƠNG 2", "Kiến trúc tổng quan — riscv_vpu_top", s4)

# ── Slide 5: Scalar Core ──────────────────────────────────────────────────────
def s5(slide):
    bullets(slide, [
        (0, "5-stage pipeline: IF → ID → EX → MEM → WB"),
        (1, "Synchronous IMEM/DMEM (1-cycle latency), Harvard architecture"),
        (1, "Full forwarding: EX→MEM, MEM→WB bypass mux"),
        (0, "Hazard detection"),
        (1, "Load-use: stall PC + IF/ID, inject bubble vào EX"),
        (1, "VPU stall: freeze toàn bộ pipeline khi vpu_ready = 0"),
        (0, "VPU Dispatch Protocol"),
        (1, "OP-V tại EX: drive vpu_insn_vld + vpu_insn + rs1/rs2"),
        (1, "vsetvl* writeback: chờ vpu_cfg_done, ghi vl vào scalar RF"),
    ], w=8.3)
    add_pic(slide, IMG + "Top-level_block_diagram.png", 8.8, 1.2, 4.3, 5.8)
wr_slide("CHƯƠNG 3", "Scalar RISC-V Core (RV32IM, 5-stage)", s5)

# ── Slide 6: VPU Overview ────────────────────────────────────────────────────
def s6(slide):
    add_txt(slide, "VPU Subsystem — vproc_system_wrapper", 1.3, 1.15, 11.0, 0.5,
            size=16, bold=True, color=NAVY)
    add_pic(slide, IMG + "vpu_top_wrapper.png", 0.5, 1.6, 12.3, 5.5)
wr_slide("CHƯƠNG 4", "VPU Architecture", s6)

# ── Slide 7: Decoder ─────────────────────────────────────────────────────────
def s7(slide):
    bullets(slide, [
        (0, "vproc_vdecoder: 32-bit instruction → 49-bit ctrl_bus (combinational)"),
        (1, "Bits [4:0] vs1_addr, [9:5] vs2_addr, [14:10] vd_addr"),
        (1, "Bits [20:15] funct6, [21] is_widen, [24] is_subtraction"),
        (1, "Bits [29] is_vector, [31] vm (insn[25]), [47] is_reduction"),
        (0, "5 nhóm lệnh chính"),
        (1, "Config (vsetvl*): 1 cycle, ST_CONFIG → update vl/vtype CSR"),
        (1, "Exec: arithmetic/logic/shift/mul/minmax → ST_EXEC"),
        (1, "Masking: compare instructions → ST_MASKING + ST_FINAL_MASKING"),
        (1, "Widening: 2-phase → ST_WIDENL + ST_WIDENH"),
        (1, "Reduction: serial accumulator → ST_REDUCTION + ST_REDUCTION_DONE"),
    ], w=8.0)
wr_slide("CHƯƠNG 4.1-4.2", "Decoder & FSM", s7)

# ── Slide 8: FSM Diagram ─────────────────────────────────────────────────────
def s8(slide):
    add_pic(slide, IMG + "fsm_diagram.png", 0.4, 1.1, 12.5, 6.1)
wr_slide("CHƯƠNG 4.2", "FSM 9-state — Execution Sequencer", s8)

# ── Slide 9: Execution Lane ───────────────────────────────────────────────────
def s9(slide):
    add_pic(slide, IMG + "executing_lane.png", 0.4, 1.1, 8.0, 5.8)
    bullets(slide, [
        (0, "4 lanes × 32-bit = 128 bits/cycle"),
        (1, "SEW slicing: 4×8b / 2×16b / 1×32b"),
        (0, "7 Functional Units"),
        (1, "Adder: VADD/VSUB/VRSUB, carry/borrow"),
        (1, "Multiplier: VMUL/VMULH/VMULHU/VMULHSU"),
        (1, "Shifter: VSLL/VSRL/VSRA (SEW-masked amount)"),
        (1, "Logic: VAND/VOR/VXOR"),
        (1, "Compare: 4-bit mask per word (vmseq,vmslt…)"),
        (1, "MinMax: signed & unsigned"),
        (1, "Reduction: serial per-element accumulator"),
    ], x=8.5, y=1.2, w=4.6)
wr_slide("CHƯƠNG 4.3", "Execution Lane (vproc_processor_lane × 4)", s9)

# ── Slide 10: VLSU Waveform ───────────────────────────────────────────────────
def s10(slide):
    add_txt(slide, "Vector LSU — Masked Store VSE8.v (vl=12, SEW=8, vm=0)", 1.3, 1.15, 11.0, 0.45,
            size=15, bold=True, color=NAVY)
    add_pic(slide, IMG + "vlsu_masked_store_wave.png", 0.4, 1.6, 12.5, 5.0)
    add_txt(slide, "3-word transfer: mem_be=4'b0101 mỗi word (elements chẵn active, lẻ masked). "
                   "Byte-enable logic: mask_be[j] = v0_flat[word_ctr×4 + j]",
            0.4, 6.65, 12.5, 0.55, size=12, color=DGRAY)
wr_slide("CHƯƠNG 4.4", "Vector LSU — VLE/VSE với masked store", s10)

# ── Slide 11: Reduction Unit ──────────────────────────────────────────────────
def s11(slide):
    add_pic(slide, IMG + "reduction_sim.png", 0.4, 1.1, 7.8, 5.8)
    bullets(slide, [
        (0, "Thiết kế ban đầu: parallel tree"),
        (1, "64 adder stages → F_max chỉ 16.44 MHz"),
        (0, "Thiết kế mới: serial per-element"),
        (1, "1 adder + register mỗi cycle"),
        (1, "N_elem cycles để tích lũy"),
        (0, "Kết quả cải thiện"),
        (1, "16.44 → 47.18 MHz (2.87×)"),
        (1, "ALM: 13,998 → 10,438 (−25%)"),
        (0, "Hỗ trợ: VREDSUM, VREDMAX,"),
        (0, "VREDMIN, VREDMAXU, VREDMINU"),
    ], x=8.6, y=1.2, w=4.5)
wr_slide("CHƯƠNG 4.5", "Reduction Unit — Serial Accumulator", s11)

# ── Slide 12: Memory & UART ───────────────────────────────────────────────────
def s12(slide):
    bullets(slide, [
        (0, "Harvard Architecture"),
        (1, "IMEM: 8 KB, read-only, 1-cycle sync, loaded via $readmemh"),
        (1, "DMEM: 64 KB, dual-port (Scalar Port A + VLSU Port B), byte-enable"),
        (0, "DMEM Memory Map"),
        (1, "0x0000–0x3FFF: R channel (Lena 128×128)"),
        (1, "0x4000–0x7FFF: G channel"),
        (1, "0x8000–0xBFFF: B channel"),
        (1, "0xC000–0xFFFF: Y channel (grayscale output)"),
        (0, "UART 8N1 (TL-UL slave @ 0xFF000000)"),
        (1, "TX/RX FIFO depth 8, 16× oversampling"),
        (1, "Registers: STAT (+0x04), TXDAT (+0x08), RXDAT (+0x0C)"),
    ])
wr_slide("CHƯƠNG 5", "Bộ nhớ & UART Peripheral", s12)

# ── Slide 13: Software ────────────────────────────────────────────────────────
def s13(slide):
    add_txt(slide, "Build Pipeline", 0.4, 1.1, 12.5, 0.45,
            size=15, bold=True, color=NAVY)
    add_pic(slide, IMG + "build_flow.png", 0.4, 1.55, 12.5, 3.1)
    bullets(slide, [
        (0, "Toolchain: riscv64-unknown-elf-gcc, -march=rv32im_zicsr_zve32x_zvl128b"),
        (0, "BT.601 VPU Kernel (lena_gray.S): 9 vector instructions per 16-pixel group"),
        (1, "vsetvli, 3×vle8.v, 3×vmulhu.vx, 2×vadd.vv, 1×vse8.v"),
        (1, "Y[i] = (R×77 + G×150 + B×29) >> 8  via vmulhu.vx SEW=8"),
    ], y=4.85)
wr_slide("CHƯƠNG 6", "Software & Firmware", s13)

# ── Slide 14: Verification ────────────────────────────────────────────────────
def s14(slide):
    bullets(slide, [
        (0, "3 mức kiểm tra (V-Model)"),
        (1, "Unit level: adder, mul, shifter, compare, reduction, VLSU"),
        (1, "Integration: tb_vproc_all_instr — 172 test cases, 13 nhóm opcode"),
        (1, "System: tb_lena_gray — 128×128 pixels, verify từng pixel vs reference"),
        (0, "Kết quả Regression"),
        (1, "172/172 PASS — ALU VV/VX/VI, VRSUB, Logic, Shift, Compare, MUL, Widening, Reduction"),
        (1, "75/75 PASS — VLSU: VLE8/16/32, VSE8/16/32, masked store"),
        (1, "16384/16384 pixels PASS — Lena grayscale, max_error = 0"),
        (0, "Transcript mô phỏng (35853 cycles)"),
    ], w=8.3)
    add_pic(slide, IMG + "lena_transcript.png", 8.8, 1.2, 4.3, 5.5)
wr_slide("CHƯƠNG 7", "Kiểm tra & Đánh giá (Verification)", s14)

# ── Slide 15: Performance ─────────────────────────────────────────────────────
def s15(slide):
    # Stat boxes
    stats = [
        ("Lena 128×128\nCycles", "35853", BLUE),
        ("Matmul 4×4\nCycles", "317", GREEN),
        ("AXPY N=16\nCycles", "221", RGBColor(0,112,192)),
        ("Speedup\nvs Scalar", "~10×", RED),
    ]
    bw = 2.9
    for i,(lbl,val,c) in enumerate(stats):
        x = 0.4 + i*3.15
        add_rect(slide, x, 1.2, bw, 2.6, c)
        add_txt(slide, val, x, 1.5, bw, 1.2, size=38, bold=True,
                color=WHITE, align=PP_ALIGN.CENTER)
        add_txt(slide, lbl, x, 2.9, bw, 0.7, size=14, color=WHITE,
                align=PP_ALIGN.CENTER)
    bullets(slide, [
        (0, "BT.601 Lena: 2.18 cycles/pixel — bandwidth-bound (1 byte/cycle = 1 element/cycle)"),
        (0, "Speedup ~10× vs scalar RV32IM cho inner loop (vl=16, SEW=8)"),
        (0, "Scalar dispatch overlap: VPU FIFO depth-8 cho phép scalar issue 4 lệnh trước khi stall"),
    ], y=4.1)
wr_slide("CHƯƠNG 7", "Performance Benchmarks", s15)

# ── Slide 16: FPGA Synthesis ──────────────────────────────────────────────────
def s16(slide):
    add_pic(slide, IMG + "Fmax.png", 0.4, 1.1, 8.0, 5.6)
    bullets(slide, [
        (0, "Fmax fast-corner: 47.18 MHz"),
        (1, "+2.247 ns slack"),
        (0, "ALM count: 10,438"),
        (0, "Fmax evolution:"),
        (1, "16.44 MHz (parallel tree)"),
        (1, "32.5 MHz (FIFO fix)"),
        (1, "47.18 MHz (serial red.)"),
        (0, "Bottleneck hiện tại:"),
        (1, "Scalar ALU carry chain"),
        (1, "~6ns slow-corner"),
    ], x=8.5, y=1.2, w=4.6)
wr_slide("CHƯƠNG 8", "FPGA Synthesis — Cyclone V", s16)

# ── Slide 17: FPGA Demo ───────────────────────────────────────────────────────
def s17(slide):
    add_txt(slide, "SW[0]=1: Ảnh RGB gốc", 1.5, 1.1, 4.5, 0.45,
            size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_txt(slide, "SW[0]=0: Grayscale BT.601", 7.3, 1.1, 4.5, 0.45,
            size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_pic(slide, IMGF + "mif_rgb.png",  1.5, 1.55, 4.5, 4.5)
    add_pic(slide, IMGF + "mif_gray.png", 7.3, 1.55, 4.5, 4.5)
    add_rect(slide, 0.0, 6.15, 13.33, 0.08, NAVY)
    add_txt(slide, "16,384/16,384 pixels verified byte-exact  |  max_error = 0  |  "
                   "Hardware VGA demo confirmed on DE10-Standard @ 40 MHz",
            0.3, 6.22, 12.7, 0.55, size=13, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
wr_slide("CHƯƠNG 8", "FPGA Demo — DE10-Standard VGA Output", s17)

# ── Slide 18: Open Issues ─────────────────────────────────────────────────────
def s18(slide):
    issues = [
        ("#8",  "HIGH",   "lsu.sv dùng async reset — phải fix để đáp ứng ASIC rule"),
        ("#9",  "HIGH",   "Chưa lint (latch audit) — cần chạy trước synthesis"),
        ("#14", "HIGH",   "CTRL_WIDTH=48 vs 49 bits — audit consumer bus width"),
        ("#13", "MEDIUM", "vslide1up/vslide1down chưa implement — cần cho FIR filter"),
        ("#15", "ASIC",   "DMEM logic array cần thay bằng SRAM macro cho tapeout"),
        ("#16", "MEDIUM", "vsetivli chưa decode (workaround: dùng AVL > 31)"),
    ]
    y = 1.2
    for (num, prio, desc) in issues:
        c = RED if prio == "HIGH" else (RGBColor(192,96,0) if prio == "ASIC" else RGBColor(0,112,192))
        add_rect(slide, 0.4, y, 0.65, 0.38, c)
        add_txt(slide, num, 0.4, y, 0.65, 0.38, size=13, bold=True,
                color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(slide, 1.1, y, 1.2, 0.38, c)
        add_txt(slide, prio, 1.1, y, 1.2, 0.38, size=11,
                color=WHITE, align=PP_ALIGN.CENTER)
        add_txt(slide, desc, 2.4, y, 10.5, 0.38, size=14, color=DGRAY)
        y += 0.45
wr_slide("CHƯƠNG 9", "Open Issues — Path to Tapeout", s18)

# ── Slide 19: Conclusion ──────────────────────────────────────────────────────
def s19(slide):
    bullets(slide, [
        (0, "Đã hoàn thành"),
        (1, "Triển khai đầy đủ Zve32x VPU — 172/172 regression PASS"),
        (1, "Tích hợp hệ thống end-to-end: scalar + VPU + UART + FPGA demo"),
        (1, "FPGA: 47.18 MHz, 10,438 ALMs — cải thiện 3× tốc độ, giảm 25% area"),
        (1, "VGA display hardware-verified trên DE10-Standard"),
        (0, "Hướng phát triển"),
        (1, "Fix async reset (#8) + lint pass (#9) → prerequisite cho ASIC"),
        (1, "Thêm vslide1up/down → mở rộng ứng dụng FIR filter"),
        (1, "Clock-gating + formal verification + STA closure @ 50 MHz"),
        (1, "Thay DMEM behavioral bằng SRAM macro → tapeout-ready"),
    ])
wr_slide("CHƯƠNG 9", "Kết luận & Hướng phát triển", s19)

wr.save(OUT_WR)
print(f"[1] Saved {OUT_WR}  ({len(wr.slides)} slides)")


# ════════════════════════════════════════════════════════════════════════════════
# PART 2 — SIMPLE PRESENTATION (content-first, plain style)
# ════════════════════════════════════════════════════════════════════════════════
p2 = Presentation()
p2.slide_width  = Inches(13.33)
p2.slide_height = Inches(7.5)
blank2 = p2.slide_layouts[6]

def ph(slide, title, items, imgs=None):
    """Plain header + bullet list slide."""
    # Header bar
    s = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.0), )
    s.fill.solid(); s.fill.fore_color.rgb = NAVY; s.line.fill.background()
    add_txt(slide, title, 0.25, 0.08, 12.8, 0.84, size=26, bold=True, color=WHITE)
    # Gray body
    s2 = slide.shapes.add_shape(1, Inches(0), Inches(1.0), Inches(13.33), Inches(6.5))
    s2.fill.solid(); s2.fill.fore_color.rgb = LGRAY; s2.line.fill.background()
    w = 8.0 if imgs else 12.8
    y = 1.3
    for (lvl, t) in items:
        p = "• " if lvl == 0 else "    – "
        c = NAVY if lvl == 0 else DGRAY
        sz = 18 if lvl == 0 else 15
        add_txt(slide, p + t, 0.4, y, w, 0.42, size=sz, bold=(lvl==0), color=c)
        y += 0.43
    if imgs:
        iy = 1.1
        iw = 4.6
        for ip in imgs:
            if os.path.exists(ip):
                slide.shapes.add_picture(ip, Inches(8.5), Inches(iy), Inches(iw), Inches(5.9/len(imgs)))
                iy += 5.9/len(imgs)

# Slide 1: Title
ts = p2.slides.add_slide(blank2)
ts.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5)).fill.solid()
ts.shapes[-1].fill.fore_color.rgb = NAVY; ts.shapes[-1].line.fill.background()
add_txt(ts, "RISC-V Vector Processing Unit", 0.5, 1.5, 12.3, 1.5, size=40, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)
add_txt(ts, "Targeting ASIC Tapeout", 0.5, 3.0, 12.3, 0.9, size=28, bold=False,
        color=CYAN, align=PP_ALIGN.CENTER)
add_txt(ts, "rv32im_zicsr_zve32x_zvl128b  |  VLEN=128  |  SEW=8/16/32  |  LMUL=1/2/4/8",
        0.5, 4.0, 12.3, 0.7, size=18, color=RGBColor(180,210,240), align=PP_ALIGN.CENTER)
add_txt(ts, "Pham Anh Khoi  ·  HCMUT  ·  Capstone Project 2  ·  June 2026",
        0.5, 5.5, 12.3, 0.6, size=16, color=RGBColor(150,190,230), align=PP_ALIGN.CENTER)

# Slide 2: System Architecture
s = p2.slides.add_slide(blank2)
ph_slide = s
s2_ = s.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
s2_.fill.solid(); s2_.fill.fore_color.rgb = LGRAY; s2_.line.fill.background()
s3_ = s.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.0))
s3_.fill.solid(); s3_.fill.fore_color.rgb = NAVY; s3_.line.fill.background()
add_txt(s, "System Architecture", 0.25, 0.08, 12.8, 0.84, size=26, bold=True, color=WHITE)
if os.path.exists(IMG + "TOP_design.png"):
    s.shapes.add_picture(IMG + "TOP_design.png", Inches(0.3), Inches(1.05), Inches(12.73), Inches(6.1))
add_txt(s, "riscv_vpu_top: scalar pipeline + VPU + IMEM/DMEM + UART on single clock domain",
        0.3, 7.1, 12.73, 0.35, size=13, color=DGRAY, align=PP_ALIGN.CENTER)

def simple_slide(title, items, imgs=None):
    slide = p2.slides.add_slide(blank2)
    ph(slide, title, items, imgs)
    return slide

simple_slide("5-Stage Scalar Core (RV32IM)", [
    (0, "Pipeline: IF → ID → EX → MEM → WB, synchronous IMEM/DMEM"),
    (1, "1-cycle read latency, Harvard architecture"),
    (0, "Full forwarding: EX→MEM, MEM→WB bypass mux"),
    (0, "Hazard detection"),
    (1, "Load-use stall: freeze PC + IF/ID, bubble into EX"),
    (1, "VPU stall: freeze all stages until vpu_ready = 1"),
    (0, "VPU dispatch: OP-V in EX stage"),
    (1, "vpu_insn_vld + vpu_insn + rs1/rs2_data to VPU"),
    (1, "vsetvl* writeback: wait cfg_done, write vl to scalar RF"),
], imgs=[IMG + "Top-level_block_diagram.png"])

simple_slide("VPU Subsystem Architecture", [
    (0, "vproc_system_wrapper: instantiates all VPU submodules"),
    (1, "Decoder (32b → 49-bit ctrl_bus), FIFO depth-8, FSM 9-state"),
    (1, "4 × vproc_vregfile (32 × 32b each), 4 processor lanes"),
    (1, "vproc_vec_lsu, vproc_reduction, vproc_vcsr"),
    (0, "49-bit ctrl_bus: vs1/vs2/vd addresses, funct6, flags (is_widen, is_masking, is_reduction…)"),
    (0, "Instruction FIFO (depth 8): decouple scalar from VPU execution"),
    (1, "Prevents duplicate push during scalar stall via vpu_ready gating"),
], imgs=[IMG + "vpu_top_wrapper.png"])

simple_slide("VPU Execution FSM — 9 States", [
    (0, "ST_IDLE (0): dispatch hub, wait for instruction from FIFO"),
    (0, "ST_CONFIG (1): vsetvl/vsetvli — 1 cycle, update CSR vl/vtype"),
    (0, "ST_EXEC (2): arithmetic/logic/shift/mul/minmax — N_elem/4 cycles"),
    (0, "ST_WIDENL/WIDENH (3/4): widening — 2-phase per group"),
    (0, "ST_MASKING (5) + ST_FINAL_MASKING (6): compare → pack mask bits"),
    (0, "ST_REDUCTION (7) + ST_REDUCTION_DONE (8): serial accumulation"),
    (1, "T_red = N_elem + 1 cycles  |  T_exec = ceil(N_elem / (32/SEW)) cycles"),
], imgs=[IMG + "fsm_diagram.png"])

simple_slide("Execution Lane & Functional Units", [
    (0, "4 lanes × 32-bit = processes 128-bit VRF register in 1 group"),
    (1, "SEW=8: 4 elements/cycle  |  SEW=16: 2/cycle  |  SEW=32: 1/cycle"),
    (0, "Adder: VADD/VSUB/VRSUB/VADC/VSBC — carry-save, no cross-lane carry"),
    (0, "Multiplier: VMUL/VMULH/VMULHU/VMULHSU, widening output lo/hi"),
    (0, "Shifter: VSLL/VSRL/VSRA — amount masked to log2(SEW) bits"),
    (0, "Compare: 4-bit result/word → mask write buffer → ST_FINAL_MASKING"),
    (0, "MinMax: signed/unsigned per-element mux"),
], imgs=[IMG + "executing_lane.png"])

simple_slide("Vector LSU — Masked Store Waveform", [
    (0, "Unit-stride: VLE8/16/32.v (load), VSE8/16/32.v (store)"),
    (0, "Prefetch pipeline: ST_IDLE issues word-0 request; data arrives on first ST_LOAD cycle"),
    (0, "Masked store (vm=0): mem_be[j] = v0_flat[word_ctr×4 + j] for SEW=8"),
    (1, "Tail suppression: last_word_be computed from vl mod 4"),
    (0, "Hazards: RAW scoreboard (vrf_busy_r[31:0]), WAW guard, WAR guard"),
], imgs=[IMG + "vlsu_masked_store_wave.png"])

simple_slide("Verification Results", [
    (0, "172/172 PASS — Instruction regression (tb_vproc_all_instr)"),
    (1, "13 sections: CONFIG, ALU VV/VX/VI, VRSUB, Logic, Shift, MinMax, Compare, MUL, Widening, LMUL=2, Reduction, OPMVX"),
    (0, "75/75 PASS — Vector LSU (tb_vproc_vlsu)"),
    (1, "VLE/VSE e8/e16/e32, random VL, masked store byte-enable"),
    (0, "16384/16384 PASS — Lena 128×128 system benchmark"),
    (1, "35853 cycles, max pixel error = 0 vs BT.601 reference"),
    (0, "Regression policy: every RTL change must pass 172-test gate"),
])

simple_slide("FPGA Synthesis & Performance", [
    (0, "Fmax: 47.18 MHz (fast-corner) — 2.87× vs original 16.44 MHz"),
    (1, "Slow-corner: −4.436 ns violation (scalar ALU carry chain bottleneck)"),
    (0, "Area: 10,438 ALMs — 25% reduction after serial reduction redesign"),
    (0, "Fmax evolution: 16.44 → 32.5 → 47.18 MHz"),
    (1, "Fix 1: serial reduction (remove 64-deep adder tree)"),
    (1, "Fix 2: break 143-node FIFO feedback loop"),
    (0, "Performance: ~10× speedup vs scalar RV32IM for BT.601 inner loop"),
    (1, "2.18 cycles/pixel, bandwidth-bound at SEW=8"),
], imgs=[IMG + "Fmax.png"])

# FPGA Demo slide
sd = p2.slides.add_slide(blank2)
bg = sd.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
bg.fill.solid(); bg.fill.fore_color.rgb = LGRAY; bg.line.fill.background()
hd = sd.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.0))
hd.fill.solid(); hd.fill.fore_color.rgb = NAVY; hd.line.fill.background()
add_txt(sd, "FPGA Demo — DE10-Standard VGA Output", 0.25, 0.08, 12.8, 0.84,
        size=26, bold=True, color=WHITE)
add_txt(sd, "RGB original (SW[0]=1)", 1.5, 1.05, 4.5, 0.45,
        size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_txt(sd, "Grayscale BT.601 (SW[0]=0)", 7.3, 1.05, 4.5, 0.45,
        size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
if os.path.exists(IMGF + "mif_rgb.png"):
    sd.shapes.add_picture(IMGF + "mif_rgb.png",  Inches(1.5), Inches(1.5), Inches(4.5), Inches(4.5))
if os.path.exists(IMGF + "mif_gray.png"):
    sd.shapes.add_picture(IMGF + "mif_gray.png", Inches(7.3), Inches(1.5), Inches(4.5), Inches(4.5))
add_txt(sd, "Both images pre-loaded via MIF  |  16384/16384 pixels byte-exact  |  max_error=0  |  Hardware verified",
        0.3, 6.25, 12.73, 0.55, size=14, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

simple_slide("Conclusion & Future Work", [
    (0, "Achievements"),
    (1, "Complete Zve32x VPU — 172/172 regression PASS"),
    (1, "End-to-end verified: UART streaming + BT.601 image processing + FPGA demo"),
    (1, "47.18 MHz, 10,438 ALMs on Cyclone V — 3× frequency, 25% area improvement"),
    (1, "VGA display hardware-verified on DE10-Standard board"),
    (0, "Path to ASIC Tapeout"),
    (1, "Fix async reset in lsu.sv (#8) — synthesis prerequisite"),
    (1, "Lint pass — latch audit (#9) required before DC/Yosys"),
    (1, "Replace behavioral DMEM with standard-cell SRAM macro (#15)"),
    (0, "Feature Extensions"),
    (1, "vslide1up/vslide1down for FIR/overlap-save workloads"),
    (1, "Formal verification + STA closure + clock-gating insertion"),
])

p2.save(OUT_PPTX)
print(f"[2] Saved {OUT_PPTX}  ({len(p2.slides)} slides)")
